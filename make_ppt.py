"""발표자료(.pptx) 생성 — 충남대 Campus ChatBot 텀프로젝트."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x00, 0x33, 0x5C)
BLUE = RGBColor(0x00, 0x6F, 0xC0)
LIGHT = RGBColor(0xEA, 0xF1, 0xF8)
GRAY = RGBColor(0x44, 0x4A, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def box(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE

    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def title_slide(title, subtitle, footer):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    rect(s, 0, 4.55, 13.333, 0.06, BLUE)
    tf = box(s, 0.9, 2.5, 11.5, 2.0)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; _set(r, 44, WHITE, True)
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = subtitle; _set(r, 22, RGBColor(0xBF, 0xD8, 0xEF))
    tf2 = box(s, 0.9, 6.4, 11.5, 0.7)
    r = tf2.paragraphs[0].add_run(); r.text = footer; _set(r, 16, RGBColor(0x9F, 0xB8, 0xD0))
    return s


def content_slide(title, bullets):
    """bullets: list of (text, level)."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 1.15, NAVY)
    rect(s, 0, 1.15, 13.333, 0.07, BLUE)
    tf = box(s, 0.6, 0.18, 12.1, 0.85)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run(); r.text = title; _set(r, 30, WHITE, True)

    body = box(s, 0.75, 1.55, 11.9, 5.6)
    first = True
    for text, lvl in bullets:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(7)
        bullet = "▪ " if lvl == 0 else "– "
        r = p.add_run()
        r.text = bullet + text
        if lvl == 0:
            _set(r, 19, NAVY, True)
        else:
            _set(r, 16, GRAY, False)
    return s


def section_slide(num, title):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 7.5, LIGHT)
    rect(s, 0, 3.0, 13.333, 1.5, NAVY)
    tf = box(s, 0.9, 3.0, 11.5, 1.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"  {num}  "; _set(r, 30, BLUE, True)
    r = p.add_run(); r.text = title; _set(r, 32, WHITE, True)
    return s


# ── 1. 표지 ──
title_slide(
    "충남대학교 Campus ChatBot",
    "질문 유형 분류 + RAG 기반 학내 정보 안내 챗봇",
    "자연어처리 텀프로젝트  |  컴퓨터융합학부  |  장윤상 (202202497)",
)

# ── 2. 개요 ──
content_slide("프로젝트 개요", [
    ("목적: 충남대 재학생을 위한 학내 정보 AI 챗봇", 0),
    ("질문 유형 분류 + RAG 기반 응답 + 실시간 정보 반영", 1),
    ("3개 Task 구성", 0),
    ("Task 1 — 질문 유형 분류기 (F1 정량 평가)", 1),
    ("Task 2 — 챗봇 모델 & Streamlit UI (정성 평가)", 1),
    ("Task 3 — 실시간 정보 반영 (tool calling, optional)", 1),
    ("5개 질문 카테고리", 0),
    ("졸업요건 · 학교 공지사항 · 학사일정 · 식단 안내 · 통학/셔틀버스", 1),
])

# ── 3. 시스템 아키텍처 ──
content_slide("전체 시스템 아키텍처", [
    ("데이터 구축 (오프라인)", 0),
    ("공개 웹 크롤링 → 개인정보 마스킹 → 청킹(300~500토큰) → bge-m3 임베딩 → ChromaDB", 1),
    ("추론 파이프라인 (질문 → 답변)", 0),
    ("① 질문 입력", 1),
    ("② LLM tool-calling 라우터(Gemma): 식단·셔틀·학사·공지 → 실시간 조회", 1),
    ("③ RAG: bge-m3 임베딩 → 벡터검색 → cross-encoder 리랭크 + 점수 컷오프", 1),
    ("④ [시스템 프롬프트 + 컨텍스트 + 질문] → Gemma 4 12B (4bit)", 1),
    ("⑤ 자연어 답변 + 출처", 1),
])

# ── 4. Task 1 ──
content_slide("Task 1 — 질문 유형 분류기", [
    ("모델: klue/bert-base를 밑바닥부터 직접 구현한 BERT에 가중치 copy", 0),
    ("임베딩 오차 < 1e-4로 구현 검증 (src/model/bert_scratch.py)", 1),
    ("백본 copy + 분류 헤드(Linear 768→5)만 학습, AdamW + warmup + AMP", 1),
    ("5 epoch / batch 32 / lr 2e-5 / max_len 128, 검증 Macro-F1 최고점만 저장", 1),
    ("학습 데이터 구축 (시드 직접 + LLM 증강)", 0),
    ("시드 train 316건 직접 작성(라벨 균등) + valid 271건 분리", 1),
    ("Claude Haiku paraphrase로 ~1만 건 증강 (말투·오타·줄임말 다양화)", 1),
    ("규정 준수: test_cls.json 학습 제외, KorQuAD 등 공개 Q&A 미사용", 0),
    ("출력: outputs/cls_output.json  →  Macro-F1 + Confusion Matrix", 1),
])

# ── 5. Task 2 ──
content_slide("Task 2 — 챗봇 모델", [
    ("Base 모델: google/gemma-4-12b-it (4bit NF4, bfloat16)", 0),
    ("파인튜닝 없이 RAG로 도메인 지식 주입 (LoRA 제거)", 1),
    ("Turing GPU 대응: bfloat16 필수 (float16 시 NaN → 출력 깨짐)", 1),
    ("멀티모달 unified 체크포인트 → AutoModelForImageTextToText 로드", 1),
    ("RAG: bge-m3 임베딩 + ChromaDB + cross-encoder 리랭크", 0),
    ("코퍼스: cnu_chunks 청크 (+ 실시간 cnu_live)", 1),
    ("컨텍스트에 없는 내용은 추측 금지 → '확인되지 않았어요'", 1),
    ("입출력: test_chat.json → chat_output.json + Streamlit UI", 0),
])

# ── 6. 데이터 구축 ──
content_slide("데이터 구축 파이프라인", [
    ("크롤링 대상 (robots.txt 준수, 3초 간격)", 0),
    ("plus.cnu / computer.cnu / job.cnu / cnucoop / mobileadmin / dorm", 1),
    ("정적: requests+BeautifulSoup, 동적: playwright, PDF: pdfplumber", 1),
    ("개인정보 보호", 0),
    ("학번·이름·연락처 등 크롤링 직후 마스킹, 원본 폐기", 1),
    ("청킹 & Q&A 생성", 0),
    ("300~500 토큰 / 50 오버랩, 문단·제목 경계 우선", 1),
    ("외부 LLM로 Q&A 자동 생성, 평가셋은 사람이 직접 검수", 1),
    ("규정 준수: KorQuAD 등 공개 Q&A 데이터셋 미사용", 0),
])

# ── 7. RAG 검색 ──
content_slide("RAG 검색 상세 (Advanced RAG)", [
    ("임베딩: BAAI/bge-m3 (다국어, 1024차원)", 0),
    ("VRAM 절약 위해 CPU 로드 → 12B LLM에 GPU 양보", 1),
    ("하이브리드 검색 + 리랭크", 0),
    ("밀집(bge-m3) + 희소(BM25) 검색 → RRF 융합", 1),
    ("cross-encoder(bge-reranker-v2-m3) 리랭크 + 점수 컷오프", 1),
    ("무관 청크는 컨텍스트뿐 아니라 출처 배지에서도 탈락 (누수 차단)", 1),
    ("부스팅 & 코퍼스", 0),
    ("학과·장학 키워드 감지 시 해당 source 우선순위 상향", 1),
    ("static(cnu_chunks) + live(cnu_live) 동시 검색, 모든 답변에 출처 표기", 1),
])

# ── 8. Task 3 ──
content_slide("Task 3 — 실시간 정보 반영", [
    ("LLM tool-calling 라우터 (src/skills/router.py)", 0),
    ("Gemma가 도구·인자를 직접 선택 (키워드 하드코딩 대체)", 1),
    ("tri-state: 도구선택 / 도구불필요 / 파싱실패 → 키워드 폴백", 1),
    ("지원 도구", 0),
    ("오늘의 식단(학생회관·기숙사) / 셔틀버스 / 학사일정 / 공지사항", 1),
    ("동작 방식", 0),
    ("cnu_live 컬렉션을 실행 시점에 갱신 후 검색", 1),
    ("도구 실패 시 RAG 폴백으로 안정성 확보", 1),
    ("입출력: test_realtime.json → realtime_output.json", 0),
])

# ── 9. 기술적 의사결정 ──
content_slide("기술적 의사결정", [
    ("Base 모델: Qwen3 → Gemma 4 12B 최종 확정", 0),
    ("멀티모달 unified 체크포인트 → AutoModelForImageTextToText", 1),
    ("Gemma 4는 torch 2.7+ 필요 → 박스 torch 2.12.1/cu126, transformers 5.12.1", 1),
    ("4bit NF4 양자화로 VRAM 절감 (Colab T4 폴백 시 RAG_MODE=naive)", 1),
    ("스킬 라우팅: 키워드 매칭 → LLM tool-calling (decision 004)", 0),
    ("검색 기본 모드: naive → rerank + 점수 컷오프 (맥락응답 정밀도↑)", 0),
    ("UI: Gradio → Streamlit (decision 003)", 0),
    ("재현성: 학습·평가 전 코드 seed=42 고정", 0),
])

# ── 10. 데모 ──
content_slide("데모 — Streamlit 챗 인터페이스", [
    ("실행: bash chatbot.sh → 배치 추론 + Streamlit UI", 0),
    ("사용자가 질문 입력 → 응답 + 출처 배지 표시", 1),
    ("멀티턴 대화 지원 (최근 대화 이력 반영)", 1),
    ("사이드바에서 검색 모드·임베딩 등 설정 조정 가능", 1),
    ("예시 질문", 0),
    ("\"컴퓨터융합학부 졸업 요건이 어떻게 되나요?\"", 1),
    ("\"오늘 학식 뭐 나와요?\"  /  \"셔틀버스 시간표 알려줘\"", 1),
    ("(시연 영상 / 스크린샷 삽입 위치)", 0),
])

# ── 11. 한계 & 마무리 ──
content_slide("한계 및 향후 과제", [
    ("한계", 0),
    ("일부 학교 서버가 데이터센터 IP(Colab) 차단 → 실시간 공지 제약", 1),
    ("실시간 크롤링은 학교 페이지 구조 변경에 민감", 1),
    ("LLM 라우터·리랭크로 질문당 생성/추론 1회 추가 → 응답 지연 증가", 1),
    ("향후 개선", 0),
    ("분류기-검색 연계(카테고리별 검색 라우팅)", 1),
    ("도구 호출 안정화 및 캐싱, 응답 지연 단축", 1),
    ("정리", 0),
    ("분류 + RAG + 실시간을 하나의 파이프라인으로 통합 구현", 1),
    ("규정 준수(데이터 출처·개인정보·평가셋 분리) 하에 동작 확인", 1),
])

# ── 12. 끝 ──
s = title_slide("감사합니다", "질문 환영합니다", "충남대학교 Campus ChatBot  |  장윤상")

out = "발표자료_장윤상.pptx"
prs.save(out)
print(f"생성 완료: {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
