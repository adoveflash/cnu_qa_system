"""랩실(지능소프트웨어 연구실) 신입생 Final Project (2) 발표자료(.pptx) 생성.

과제 명세(신입생_튜토리얼 'Final Project (2)')가 요구하는 발표 6항목을 그대로 골격으로 삼는다:
  ① 데이터 선정 및 분석  ② 데이터 전처리  ③ 모델 선정 기준(왜 이 모델?)
  ④ 구현 방식 설명       ⑤ 정성 평가       ⑥ 향후 연구
제약: 오픈 웨이트 모델만(GPT·Claude 금지) / 기존 데이터셋 금지·직접 구축 / Vector DB 자유 / Streamlit.
사용자가 고민한 '키워드→LLM tool-calling 라우팅' 회고를 ④·⑥에 배치.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x00, 0x33, 0x5C)
BLUE = RGBColor(0x00, 0x6F, 0xC0)
LIGHT = RGBColor(0xEA, 0xF1, 0xF8)
GRAY = RGBColor(0x44, 0x4A, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x7D, 0x4F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def box(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame


def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def _header(s, title):
    rect(s, 0, 0, 13.333, 1.15, NAVY)
    rect(s, 0, 1.15, 13.333, 0.07, BLUE)
    tf = box(s, 0.6, 0.18, 12.1, 0.85)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 28, WHITE, True)


def title_slide(title, subtitle, footer):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    rect(s, 0, 4.55, 13.333, 0.06, BLUE)
    tf = box(s, 0.9, 2.4, 11.5, 2.2)
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 40, WHITE, True)
    p2 = tf.add_paragraph()
    r = p2.add_run()
    r.text = subtitle
    _set(r, 20, RGBColor(0xBF, 0xD8, 0xEF))
    tf2 = box(s, 0.9, 6.4, 11.5, 0.7)
    r = tf2.paragraphs[0].add_run()
    r.text = footer
    _set(r, 15, RGBColor(0x9F, 0xB8, 0xD0))
    return s


def content_slide(title, bullets):
    """bullets: list of (text, level)."""
    s = prs.slides.add_slide(BLANK)
    _header(s, title)
    body = box(s, 0.75, 1.55, 11.9, 5.6)
    first = True
    for text, lvl in bullets:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(8)
        bullet = "▪ " if lvl == 0 else "– "
        r = p.add_run()
        r.text = bullet + text
        if lvl == 0:
            _set(r, 19, NAVY, True)
        else:
            _set(r, 16, GRAY, False)
    return s


def trouble_slide(title, intro, rows, heads=("증상", "원인", "해결")):
    """3열 표 슬라이드. rows = [(a, b, c)]."""
    s = prs.slides.add_slide(BLANK)
    _header(s, title)
    if intro:
        tf = box(s, 0.75, 1.35, 11.9, 0.6)
        r = tf.paragraphs[0].add_run()
        r.text = intro
        _set(r, 15, GRAY, False)
    n = len(rows) + 1
    top = 2.05
    height = min(4.9, 0.9 * n)
    table = s.shapes.add_table(n, 3, Inches(0.6), Inches(top), Inches(12.1), Inches(height)).table
    table.columns[0].width = Inches(3.4)
    table.columns[1].width = Inches(4.35)
    table.columns[2].width = Inches(4.35)
    for j, h in enumerate(heads):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = h
        _set(r, 15, WHITE, True)
    for i, row in enumerate(rows, start=1):
        cols = [(row[0], RED), (row[1], GRAY), (row[2], GREEN)]
        for j, (txt, col) in enumerate(cols):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = txt
            _set(r, 12.5, col, j == 0)
    return s


# ── 1. 표지 ──
title_slide(
    "충남대 학내 정보 RAG Q/A 챗봇",
    "지능소프트웨어 연구실 신입생 Final Project (2) — RAG 기반 Q/A System",
    "지능소프트웨어 연구실  |  컴퓨터융합학부  |  장윤상 (202202497)",
)

# ── 2. 목표 & 요구사항 ──
content_slide("프로젝트 목표 & 과제 요구사항", [
    ("Final Project (2): RAG를 사용한 Q/A System 구축", 0),
    ("Chat Interface(Streamlit)를 통한 정성 평가 — working chatbot + API", 1),
    ("도메인: 충남대 재학생 대상 학내 정보 QA", 0),
    ("졸업요건 · 공지 · 학사일정 · 식단 · 셔틀 등 5개 영역", 1),
    ("과제 제약 준수", 0),
    ("오픈 웨이트 모델만 사용 (GPT·Claude 등 closed 모델 금지)", 1),
    ("기존 데이터셋 사용 금지 → 직접 수집·구축", 1),
    ("Vector DB 자유 선택(ChromaDB) / 데모는 Streamlit 필수", 1),
])

# ── 3. ① 데이터 선정 및 분석 ──
content_slide("① 데이터 선정 및 분석", [
    ("왜 직접 구축했나", 0),
    ("학내 정보는 공개 QA 데이터셋에 없음 + 과제 규정상 기존 데이터셋 금지", 1),
    ("KorQuAD 등 공개 Q&A 데이터셋 미사용 (규정 준수)", 1),
    ("수집 대상 — 충남대 공개 웹 (robots.txt 준수, 요청 간 3초)", 0),
    ("plus.cnu / computer.cnu / job.cnu / cnucoop / mobileadmin / dorm", 1),
    ("정적: requests+BeautifulSoup, 동적: playwright, PDF: pdfplumber", 1),
    ("5개 도메인을 균형 있게 수집해 카테고리 편향 최소화", 0),
])

# ── 4. ② 데이터 전처리 ──
content_slide("② 데이터 전처리 → 벡터 색인", [
    ("정제 & 개인정보 보호", 0),
    ("HTML/PDF 본문 추출, 메뉴·중복·잡음 제거 (cleaned)", 1),
    ("학번·이름·연락처 등 크롤링 직후 마스킹, 원본 폐기", 1),
    ("청킹", 0),
    ("300~500 토큰 / 50 토큰 오버랩, 문단·제목 경계 우선", 1),
    ("임베딩 & 색인 (RAG 코퍼스)", 0),
    ("정제 청크(chunks.jsonl)를 bge-m3로 임베딩 → ChromaDB 적재", 1),
    ("※ 파인튜닝 없이 base 모델 직접 사용 → 검색 코퍼스는 '문서 청크'", 1),
    ("Q&A 쌍은 정성·정량 평가용으로만 별도 구축(추론엔 미사용)", 1),
])

# ── 5. ③ 모델 선정 기준 ──
content_slide("③ 모델 선정 기준 — 왜 이 모델?", [
    ("생성 LLM: google/gemma-4-12b-it (오픈 웨이트)", 0),
    ("closed 금지 규정 충족 + 한국어·지시 이행 성능 우수", 1),
    ("4bit NF4 양자화 → 12B를 단일 GPU/Colab T4에서 구동", 1),
    ("파인튜닝 대신 RAG로 최신·정확 지식 주입 (환각 억제)", 1),
    ("임베딩: BAAI/bge-m3 (다국어·한국어, 1024차원)", 0),
    ("CPU 로드로 VRAM을 12B LLM에 양보", 1),
    ("Vector DB: ChromaDB (로컬 구동, 제출물에 번들)", 0),
])

# ── 6. ④ 구현 — 시스템 아키텍처 ──
content_slide("④ 구현 — 시스템 아키텍처", [
    ("과제 표준 파이프라인 준수", 0),
    ("Documents → Embedding Model → Vector DB → LLM → Answer", 1),
    ("User ↔ Chat Bot Web App (Streamlit)", 1),
    ("추론 흐름 (질문 → 답변)", 0),
    ("① 라우터: 실시간 스킬 필요 여부 판단 (식단·셔틀·학사·공지)", 1),
    ("② RAG 검색 → 관련 청크로 컨텍스트 구성", 1),
    ("③ [시스템 프롬프트 + 컨텍스트 + 질문] → Gemma → 답변 + 출처", 1),
])

# ── 7. ④ 구현 — RAG 검색 설계 ──
content_slide("④ 구현 — RAG 검색 설계 (맥락 정밀도)", [
    ("문제: naive top-k 검색은 무관 청크가 컨텍스트·출처로 누수", 0),
    ("엉뚱한 학과·부서 출처가 답변 근거로 붙던 현상", 1),
    ("해결: Advanced RAG 파이프라인", 0),
    ("밀집(bge-m3) + 희소(BM25) 하이브리드 검색 → RRF 융합", 1),
    ("cross-encoder(bge-reranker-v2-m3)로 재점수화 (static/live 동일 스케일)", 1),
    ("점수 컷오프: 관련도 낮은 청크는 컨텍스트·출처 배지에서 모두 탈락", 1),
    ("전부 컷오프되면 빈 컨텍스트 → '확인되지 않았어요' 정직 응답", 1),
])

# ── 8. ④ 구현 — 실시간 라우팅 회고 (문제점) ──
trouble_slide(
    "④ 구현 — 실시간 라우팅의 재설계 (회고)",
    "처음엔 키워드 매칭으로 스킬을 라우팅 → 실제 구동·테스트에서 4가지 한계가 드러남",
    [
        ("정적 정보까지 매번 실시간 조회",
         "셔틀·1학 메뉴는 학기 단위 고정인데 매 질문 서버 조회 (느림·불안정)",
         "정적/동적 미분리 → '메뉴를 상수로' 땜질 커밋까지 발생"),
        ("변형·줄임말 인식 불가",
         "'1학'(제1학생회관)을 추론 못 해 엉뚱한 답변",
         "'1학'·'일학'·'1학관'을 사람이 무한 하드코딩해야 함"),
        ("의미를 고려하지 않는 매칭",
         "'의미 이해'를 LLM이 아니라 사람이 키워드로 대신",
         "일반화 불가 — 새 표현마다 규칙 추가"),
        ("fallback이 오답을 부름",
         "실시간 조회 실패 → 과거 RAG 데이터로 잘못된 정보 전달 위험",
         "차라리 '최신 확인 실패'라 답하는 게 정직"),
    ],
    heads=("한계", "드러난 상황", "함의 / 결론"),
)

# ── 9. ④ 구현 — 키워드 → LLM tool-calling ──
content_slide("④ 구현 — 키워드 → LLM tool-calling 전환", [
    ("근본 원인: 의미 판단을 LLM 아닌 키워드 규칙(사람)이 대신한 것", 0),
    ("→ 빠름·재현의 대가로 일반화 불가 + 하드코딩 + 정적/동적 실패", 1),
    ("개선: Gemma가 도구·인자를 직접 선택 (src/skills/router.py)", 0),
    ("도구 스펙(JSON 스키마) 제공 → 도구 선택 + 인자 추출을 한 번에", 1),
    ("'1학'→제1학생회관, '내일'→날짜 등 의미 기반 인자 추론", 1),
    ("tri-state: 도구선택 / 도구불필요 / 파싱실패 → 키워드 폴백(안전망)", 1),
    ("(구동 트러블슈팅) cu126 torch, unified 체크포인트 로드, bf16 dtype", 0),
])

# ── 10. ⑤ 정성 평가 & 데모 ──
content_slide("⑤ 정성 평가 & 데모", [
    ("실행: bash chatbot.sh → 배치 추론 + Streamlit UI", 0),
    ("텍스트 입력창 + 응답 + 출처 배지 + 멀티턴 대화", 1),
    ("정성 개선 관측", 0),
    ("키워드 → LLM 라우팅: '1학' 등 변형 질문에서 도구 호출 정확도↑", 1),
    ("naive → rerank+컷오프: 무관 출처 누수 차단, 근거 정합성↑", 1),
    ("예시 질문", 0),
    ("\"컴퓨터융합학부 졸업 요건?\"  /  \"오늘 1학 뭐 나와?\"  /  \"셔틀버스 시간표\"", 1),
    ("(라이브 데모 / 시연 스크린샷 삽입 위치)", 0),
])

# ── 11. ⑥ 향후 연구 ──
content_slide("⑥ 향후 연구", [
    ("라우팅 고도화", 0),
    ("임베딩 유사도 라우팅 vs LLM 라우팅 — 비용 대비 정확도 비교", 1),
    ("정적/동적 정보 분리 원칙화 (셔틀·상수 vs 진짜 실시간 tool)", 1),
    ("신뢰성", 0),
    ("fallback 투명화 — 최신 조회 실패 시 정직 명시 + 데이터 날짜 표기", 1),
    ("리랭커 컷오프 임계값 캘리브레이션(bge-reranker logit 분포)", 1),
    ("성능", 0),
    ("라우터 JSON 출력 안정화 + 응답 지연 단축(캐싱), 추가 LLM 훈련 실험", 1),
])

# ── 12. 마무리 ──
title_slide("감사합니다", "Q & A", "지능소프트웨어 연구실  |  충남대 컴퓨터융합학부  |  장윤상")

out = "발표자료_랩실_장윤상.pptx"
prs.save(out)
print(f"생성 완료: {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
